"""
Advanced Markdown -> PPTX converter (overflow-safe, images + tables supported).

Features:
- Automatic slide splitting to avoid overflow (PDF-safe)
- **Bold text preserved** using PowerPoint runs
- Markdown images ![](Images/xxx.png) automatically inserted IN ORDER
- Markdown tables converted into real PPTX tables
- Images assumed to live in ./Images relative to script

Markdown handling:
- # / ## / ### headings create slide breaks
- Bullet lists preserved
- Tables (| --- |) become PPTX tables
- Images: inserted as full-width image slides

Dependencies:
    python -m pip install python-pptx

Usage:
    python markdown_to_pptx.py input.md output.pptx
"""

import sys
import re
import os
from pptx import Presentation
from pptx.util import Pt, Inches

MAX_LINES_PER_SLIDE = 12  # conservative
IMAGE_DIR = "Images"


def parse_bold_runs(text):
    parts = re.split(r"(\*\*.*?\*\*)", text)
    runs = []
    for part in parts:
        if part.startswith('**') and part.endswith('**'):
            runs.append((part[2:-2], True))
        elif part:
            runs.append((part, False))
    return runs


def clean_md(line):
    return re.sub(r"[`*_]", "", line).strip()


def chunk_content(lines, max_lines):
    for i in range(0, len(lines), max_lines):
        yield lines[i:i + max_lines]


def is_image(line):
    return re.match(r"!\[.*?\]\((.*?)\)", line)


def is_table_sep(line):
    return re.match(r"\|?\s*-+", line)


def markdown_to_pptx(md_path, pptx_path):
    with open(md_path, encoding="utf-8") as f:
        lines = f.readlines()

    prs = Presentation()

    current_title = None
    buffer = []
    table_buffer = []

    def flush_buffer():
        nonlocal buffer
        if not buffer:
            return
        for chunk in chunk_content(buffer, MAX_LINES_PER_SLIDE):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = current_title or ""
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for i, line in enumerate(chunk):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                p.clear()
                for text, is_bold in parse_bold_runs(line):
                    run = p.add_run()
                    run.text = clean_md(text)
                    run.font.bold = is_bold
                    run.font.size = Pt(18)
                p.level = 1 if line.strip().startswith(('-', '*')) else 0
        buffer = []

    def flush_table():
        nonlocal table_buffer
        if not table_buffer:
            return
        rows = [
            [clean_md(cell) for cell in row.strip('|').split('|')]
            for row in table_buffer
            if not is_table_sep(row)
        ]
        if not rows:
            table_buffer = []
            return
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = current_title or "Table"
        rows_n = len(rows)
        cols_n = len(rows[0])
        table = slide.shapes.add_table(
            rows_n,
            cols_n,
            Inches(0.5),
            Inches(1.5),
            Inches(9),
            Inches(5)
        ).table
        for r in range(rows_n):
            for c in range(cols_n):
                cell_tf = table.cell(r, c).text_frame
                cell_tf.text = rows[r][c]
        table_buffer = []

    for raw in lines:
        line = raw.rstrip()

        img_match = is_image(line)
        if img_match:
            flush_buffer()
            flush_table()
            img_path = img_match.group(1)
            img_file = os.path.join(IMAGE_DIR, os.path.basename(img_path))
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = current_title or ""
            if os.path.exists(img_file):
                slide.shapes.add_picture(img_file, Inches(1), Inches(1.5), width=Inches(8))
            else:
                tx = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
                tx.text_frame.text = f"[Missing image: {img_file}]"
            continue

        if line.startswith('# '):
            flush_buffer()
            flush_table()
            current_title = clean_md(line[2:])
            continue

        if line.startswith('## ') or line.startswith('### '):
            flush_buffer()
            flush_table()
            current_title = clean_md(line.split(' ', 1)[1])
            continue

        if line.strip().startswith('|'):
            table_buffer.append(line)
            continue

        if table_buffer:
            flush_table()

        if line.strip():
            buffer.append(line)

    flush_buffer()
    flush_table()
    prs.save(pptx_path)


if __name__ == '__main__':
    if len(sys.argv) != 3:
        print("Usage: python markdown_to_pptx.py input.md output.pptx")
        sys.exit(1)

    markdown_to_pptx(sys.argv[1], sys.argv[2])
